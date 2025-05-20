import streamlit as st
import os
import tempfile
from pptx import Presentation
import difflib
from io import BytesIO
import hashlib
from PIL import Image
import io
import base64

st.set_page_config(
    page_title="PowerPoint Comparison Tool",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="collapsed"
)

# Custom CSS for better appearance
st.markdown("""
<style>
    .main {
        padding: 2rem;
    }
    .stButton>button {
        width: 100%;
    }
    .diff-added {
        background-color: #d4edda;
        color: #155724;
        padding: 2px 4px;
        border-radius: 3px;
    }
    .diff-removed {
        background-color: #f8d7da;
        color: #721c24;
        padding: 2px 4px;
        border-radius: 3px;
    }
    .header-container {
        display: flex;
        align-items: center;
        gap: 10px;
    }
    .result-container {
        margin-top: 20px;
        padding: 20px;
        border-radius: 5px;
        background-color: #f8f9fa;
    }
</style>
""", unsafe_allow_html=True)

def extract_images_from_shape(shape, parent_info=""):
    """Recursively extract image data from a shape and any child shapes."""
    images = []
    shape_type = getattr(shape, "shape_type", None)
    debug_info = f"{parent_info} -> Shape type: {shape_type}"
    
    try:
        # Method 1: Direct image extraction for picture shapes
        if shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            try:
                image_blob = shape.image.blob
                img_data = process_image_blob(image_blob, debug_info)
                if img_data:
                    images.append(img_data)
            except Exception as e:
                pass  # Try alternative methods
        
        # Method 2: Check for image in shape fill
        try:
            if hasattr(shape, "fill") and shape.fill.type == 2:  # MSO_FILL.PICTURE
                image_blob = shape.fill.fore_color.picture.blob
                img_data = process_image_blob(image_blob, debug_info + " (from fill)")
                if img_data:
                    images.append(img_data)
        except Exception as e:
            pass
            
        # Method 3: Check placeholders that might contain images
        if shape_type == 14:  # MSO_SHAPE_TYPE.PLACEHOLDER
            try:
                # Some placeholders can contain images
                if hasattr(shape, "image") and hasattr(shape.image, "blob"):
                    image_blob = shape.image.blob
                    img_data = process_image_blob(image_blob, debug_info + " (from placeholder)")
                    if img_data:
                        images.append(img_data)
            except Exception as e:
                pass
                
        # Method 4: Recursively process group shapes
        if shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            try:
                for child_shape in shape.shapes:
                    child_images = extract_images_from_shape(child_shape, debug_info + " (group child)")
                    images.extend(child_images)
            except Exception as e:
                pass
    except Exception as e:
        pass
        
    return images

def process_image_blob(image_blob, debug_info=""):
    """Process an image blob to extract metadata and create thumbnail."""
    try:
        image_hash = hashlib.md5(image_blob).hexdigest()
        
        # Get image dimensions
        try:
            img = Image.open(io.BytesIO(image_blob))
            width, height = img.size
            image_format = img.format
        except:
            width, height = 0, 0
            image_format = "Unknown"
        
        # Create base64 thumbnail for display
        try:
            img = Image.open(io.BytesIO(image_blob))
            img.thumbnail((100, 100))
            buffered = io.BytesIO()
            img.save(buffered, format="PNG")
            img_str = base64.b64encode(buffered.getvalue()).decode()
            thumbnail = f"data:image/png;base64,{img_str}"
        except:
            thumbnail = None
            
        return {
            "hash": image_hash,
            "size": len(image_blob),
            "width": width,
            "height": height,
            "format": image_format,
            "thumbnail": thumbnail,
            "debug_info": debug_info
        }
    except:
        return None

def extract_content_from_pptx(file_content):
    """Extract all text and image content from a PowerPoint file."""
    presentation = Presentation(BytesIO(file_content))
    all_slides_content = []
    
    for slide_number, slide in enumerate(presentation.slides, 1):
        slide_text = []
        slide_images = []
        
        # Extract content from all shapes
        for shape in slide.shapes:
            # Extract text
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
                
            # Handle tables
            if hasattr(shape, "has_table") and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            slide_text.append(cell.text.strip())
            
            # Extract images recursively from all possible locations
            images = extract_images_from_shape(shape, f"Slide {slide_number}")
            if images:
                slide_images.extend(images)
        
        # Join all text from this slide
        all_slides_content.append({
            "slide_number": slide_number,
            "text_content": "\n".join(slide_text),
            "images": slide_images,
            "image_count": len(slide_images)
        })
    
    return all_slides_content

def format_diff_line(line):
    """Format diff lines with color highlighting."""
    if line.startswith('+ '):
        return f'<span class="diff-added">{line}</span>'
    elif line.startswith('- '):
        return f'<span class="diff-removed">{line}</span>'
    else:
        return line

def compare_presentations(file1_content, file2_content):
    """Compare text and image content between two PowerPoint files."""
    try:
        slides1 = extract_content_from_pptx(file1_content)
        slides2 = extract_content_from_pptx(file2_content)
        
        # Check if slide counts are different but continue with comparison
        slide_count_diff = len(slides1) != len(slides2)
        if slide_count_diff:
            max_slides = max(len(slides1), len(slides2))
            min_slides = min(len(slides1), len(slides2))
            slide_count_message = f"Presentations have different number of slides: {len(slides1)} vs {len(slides2)}. Comparing the first {min_slides} slides."
        else:
            slide_count_message = ""
        
        # Compare each slide up to the minimum number of slides in both presentations
        differences = []
        identical_count = 0
        text_diff_count = 0
        image_diff_count = 0
        min_slides = min(len(slides1), len(slides2))
        
        for i in range(min_slides):
            slide1 = slides1[i]
            slide2 = slides2[i]
            slide_num = slide1["slide_number"]
            content1 = slide1["text_content"]
            content2 = slide2["text_content"]
            images1 = slide1["images"]
            images2 = slide2["images"]
            
            slide_differences = {
                "slide_number": slide_num,
                "has_text_diff": False,
                "has_image_diff": False,
                "image_count1": len(images1),
                "image_count2": len(images2)
            }
            
            # Check text differences
            if content1 != content2:
                # Generate diff
                diff = list(difflib.ndiff(content1.splitlines(), content2.splitlines()))
                
                slide_differences["has_text_diff"] = True
                slide_differences["text_diff"] = diff
                slide_differences["ppt1_content"] = content1
                slide_differences["ppt2_content"] = content2
                text_diff_count += 1
            
            # Check image differences
            image_hashes1 = {img["hash"]: img for img in images1}
            image_hashes2 = {img["hash"]: img for img in images2}
            
            # Images in first PPT but not in second
            missing_in_ppt2 = []
            for h, img in image_hashes1.items():
                if h not in image_hashes2:
                    missing_in_ppt2.append(img)
            
            # Images in second PPT but not in first
            missing_in_ppt1 = []
            for h, img in image_hashes2.items():
                if h not in image_hashes1:
                    missing_in_ppt1.append(img)
            
            # Simple image count check - helps catch cases where extraction might miss some images
            if len(images1) != len(images2):
                slide_differences["image_count_different"] = True
            
            if missing_in_ppt1 or missing_in_ppt2 or slide_differences.get("image_count_different", False):
                slide_differences["has_image_diff"] = True
                slide_differences["images_missing_in_ppt1"] = missing_in_ppt1
                slide_differences["images_missing_in_ppt2"] = missing_in_ppt2
                image_diff_count += 1
            
            # Check if this slide has any differences
            if slide_differences["has_text_diff"] or slide_differences["has_image_diff"]:
                differences.append(slide_differences)
            else:
                identical_count += 1
        
        # Add information about extra slides
        extra_slides = []
        if len(slides1) > len(slides2):
            for i in range(min_slides, len(slides1)):
                extra_slides.append({
                    "slide_number": slides1[i]["slide_number"],
                    "in_presentation": 1,
                    "text_content": slides1[i]["text_content"],
                    "image_count": len(slides1[i]["images"])
                })
        elif len(slides2) > len(slides1):
            for i in range(min_slides, len(slides2)):
                extra_slides.append({
                    "slide_number": slides2[i]["slide_number"],
                    "in_presentation": 2,
                    "text_content": slides2[i]["text_content"],
                    "image_count": len(slides2[i]["images"])
                })
        
        # Generate result
        if not differences and not slide_count_diff:
            return {
                "identical": True,
                "summary": f"All {len(slides1)} slides have identical text and image content",
                "differences": []
            }
        else:
            total_compared = min_slides
            return {
                "identical": False,
                "summary": f"Found differences in {len(differences)} out of {total_compared} compared slides ({identical_count} slides identical)",
                "detailed_summary": f"Text differences: {text_diff_count} slides, Image differences: {image_diff_count} slides",
                "differences": differences,
                "extra_slides": extra_slides,
                "slide_count_different": slide_count_diff,
                "slide_count_message": slide_count_message if slide_count_diff else ""
            }
            
    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        return {
            "identical": False,
            "summary": f"Error comparing presentations: {str(e)}",
            "error_details": error_details,
            "differences": [],
            "error": True
        }

# App layout with debug mode toggle
st.title("📊 PowerPoint Comparison Tool")
st.markdown("""
This tool compares both text and image content between two PowerPoint presentations and identifies differences.
Upload two .pptx files to get started.
""")

# Add a debug mode in the sidebar
with st.sidebar:
    st.title("Settings")
    debug_mode = st.checkbox("Debug Mode", value=False, help="Show additional debugging information")

# File uploaders
col1, col2 = st.columns(2)
with col1:
    st.subheader("First Presentation")
    file1 = st.file_uploader("Upload first PowerPoint file", type=["pptx"], key="file1")
    if file1:
        st.success(f"Uploaded: {file1.name}")

with col2:
    st.subheader("Second Presentation")
    file2 = st.file_uploader("Upload second PowerPoint file", type=["pptx"], key="file2")
    if file2:
        st.success(f"Uploaded: {file2.name}")

# Compare button
if file1 and file2:
    if st.button("Compare Presentations", type="primary"):
        with st.spinner("Comparing presentations..."):
            result = compare_presentations(file1.getvalue(), file2.getvalue())
            
            # Display summary
            if result.get("error", False):
                st.error(result["summary"])
                if "error_details" in result:
                    with st.expander("Error Details"):
                        st.code(result["error_details"])
            elif result["identical"]:
                st.success(result["summary"])
            else:
                st.warning(result["summary"])
                
                # If different slide count
                if result.get("slide_count_different", False):
                    st.info(result["slide_count_message"])
                
                # Display detailed summary
                if "detailed_summary" in result:
                    st.info(result["detailed_summary"])
                
                # Display differences
                st.subheader("Detailed Differences")
                
                # Display extra slides if any
                if "extra_slides" in result and result["extra_slides"]:
                    st.subheader("Extra Slides")
                    for slide in result["extra_slides"]:
                        pres_num = slide["in_presentation"]
                        slide_num = slide["slide_number"]
                        with st.expander(f"Slide {slide_num} (only in Presentation {pres_num})"):
                            st.markdown(f"**Text Content:**")
                            st.text(slide["text_content"])
                            st.markdown(f"**Image Count:** {slide['image_count']}")
                
                # Display differences in common slides
                if result["differences"]:
                    st.subheader("Differences in Common Slides")
                    
                    for diff_data in result["differences"]:
                        slide_num = diff_data["slide_number"]
                        diff_types = []
                        
                        if diff_data.get("has_text_diff", False):
                            diff_types.append("Text")
                        if diff_data.get("has_image_diff", False):
                            diff_types.append("Images")
                        
                        diff_type_str = " & ".join(diff_types)
                        
                        with st.expander(f"Slide {slide_num} - {diff_type_str} Difference"):
                            # Text differences
                            if diff_data.get("has_text_diff", False):
                                st.markdown("### Text Differences")
                                col1, col2 = st.columns(2)
                                
                                with col1:
                                    st.markdown("**First Presentation**")
                                    st.text(diff_data["ppt1_content"])
                                    
                                with col2:
                                    st.markdown("**Second Presentation**")
                                    st.text(diff_data["ppt2_content"])
                                
                                st.markdown("**Detailed Text Differences**")
                                formatted_diff = "<br>".join([format_diff_line(line) for line in diff_data["text_diff"]])
                                st.markdown(f"<div style='background-color: #f8f9fa; padding: 10px; border-radius: 5px;'>{formatted_diff}</div>", unsafe_allow_html=True)
                            
                            # Image differences
                            if diff_data.get("has_image_diff", False):
                                st.markdown("### Image Differences")
                                
                                # Show image counts
                                img_count1 = diff_data.get("image_count1", 0)
                                img_count2 = diff_data.get("image_count2", 0)
                                st.markdown(f"**Image count:** First presentation: {img_count1}, Second presentation: {img_count2}")
                                
                                missing_in_ppt1 = diff_data.get("images_missing_in_ppt1", [])
                                missing_in_ppt2 = diff_data.get("images_missing_in_ppt2", [])
                                
                                col1, col2 = st.columns(2)
                                
                                with col1:
                                    st.markdown("**Images only in First Presentation**")
                                    if missing_in_ppt2:
                                        for i, img in enumerate(missing_in_ppt2):
                                            if img.get("thumbnail"):
                                                st.markdown(f"""
                                                <div style="border: 1px solid #ddd; padding: 10px; margin-bottom: 10px; border-radius: 5px;">
                                                    <img src="{img['thumbnail']}" style="max-width: 100px; max-height: 100px;">
                                                    <p>Size: {img['width']}x{img['height']} px, Format: {img['format']}</p>
                                                </div>
                                                """, unsafe_allow_html=True)
                                            else:
                                                st.markdown(f"Image {i+1}: Unable to display preview")
                                    else:
                                        st.markdown("*No unique images*")
                                
                                with col2:
                                    st.markdown("**Images only in Second Presentation**")
                                    if missing_in_ppt1:
                                        for i, img in enumerate(missing_in_ppt1):
                                            if img.get("thumbnail"):
                                                st.markdown(f"""
                                                <div style="border: 1px solid #ddd; padding: 10px; margin-bottom: 10px; border-radius: 5px;">
                                                    <img src="{img['thumbnail']}" style="max-width: 100px; max-height: 100px;">
                                                    <p>Size: {img['width']}x{img['height']} px, Format: {img['format']}</p>
                                                </div>
                                                """, unsafe_allow_html=True)
                                            else:
                                                st.markdown(f"Image {i+1}: Unable to display preview")
                                    else:
                                        st.markdown("*No unique images*")

# Additional information
with st.expander("How to use this tool"):
    st.markdown("""
    ### Instructions
    1. Upload your first PowerPoint presentation using the uploader on the left
    2. Upload your second PowerPoint presentation using the uploader on the right
    3. Click the "Compare Presentations" button
    4. Review the results:
       - Green success message means the presentations are identical
       - Warning message means differences were found
       - Expand each slide section to see detailed differences
    
    ### Reading the differences
    - Text with a green background (prefixed with `+ `) is present in the second presentation but not in the first
    - Text with a red background (prefixed with `- `) is present in the first presentation but not in the second
    - Unchanged text is shown without color highlighting
    
    ### Limitations
    - Complex PowerPoint elements like SmartArt may not be fully analyzed
    - Very large presentations may take longer to process
    - Image comparison is based on content hash - images that look similar but have different encodings will be detected as different
    - Some image types may not display properly in the comparison thumbnails
    """)

with st.expander("About this tool"):
    st.markdown("""
    This PowerPoint Text and Image Comparison Tool was built with:
    - Python
    - Streamlit for the web interface
    - python-pptx for PowerPoint file processing
    - difflib for text comparison
    - PIL (Pillow) for image processing
    
    The tool compares both text content and images between PowerPoint presentations.
    """)